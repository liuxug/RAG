import fitz
import markdown
from pathlib import Path
from loguru import logger
from typing import List, Tuple
from docx import Document as DocxDocument
from pptx import Presentation
from openpyxl import load_workbook
from PIL import Image
from odf.opendocument import load
from odf.table import Table, TableRow, TableCell
from odf.text import P
import io


class DocumentProcessor:
    SUPPORTED_EXTENSIONS = {
        ".pdf", ".md", ".txt", ".docx", ".doc",
        ".xlsx", ".xls", ".pptx", ".ppt",
        ".jpg", ".jpeg", ".png", ".gif", ".bmp",
        ".csv", ".json", ".xml", ".html", ".htm",
        ".ods", ".odt"
    }

    @staticmethod
    def _extract_text_from_pdf(file_path: Path) -> List[Tuple[int, str]]:
        """Extract text from PDF file, returning list of (page_number, text)"""
        pages = []
        try:
            doc = fitz.open(file_path)
            for page_num in range(len(doc)):
                page = doc[page_num]
                text = page.get_text("text")
                if text.strip():
                    pages.append((page_num + 1, text.strip()))
            doc.close()
            logger.info(f"Extracted {len(pages)} pages from PDF: {file_path.name}")
        except Exception as e:
            logger.error(f"Failed to extract text from PDF: {file_path.name}, error: {str(e)}")
            raise
        return pages

    @staticmethod
    def _extract_text_from_markdown(file_path: Path) -> List[Tuple[int, str]]:
        """Extract text from Markdown file"""
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                content = f.read()
            html = markdown.markdown(content)
            text = "".join([p.strip() for p in html.replace("<p>", "\n").replace("</p>", "\n").split("\n") if p.strip()])
            logger.info(f"Extracted text from Markdown: {file_path.name}")
            return [(1, text)]
        except Exception as e:
            logger.error(f"Failed to extract text from Markdown: {file_path.name}, error: {str(e)}")
            raise

    @staticmethod
    def _extract_text_from_txt(file_path: Path) -> List[Tuple[int, str]]:
        """Extract text from TXT file"""
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                content = f.read()
            logger.info(f"Extracted text from TXT: {file_path.name}")
            return [(1, content)]
        except Exception as e:
            logger.error(f"Failed to extract text from TXT: {file_path.name}, error: {str(e)}")
            raise

    @staticmethod
    def _extract_text_from_docx(file_path: Path) -> List[Tuple[int, str]]:
        """Extract text from DOCX file"""
        try:
            doc = DocxDocument(file_path)
            paragraphs = []
            for para in doc.paragraphs:
                if para.text.strip():
                    paragraphs.append(para.text.strip())
            
            for table in doc.tables:
                for row in table.rows:
                    row_text = "\t".join(cell.text.strip() for cell in row.cells)
                    if row_text.strip():
                        paragraphs.append(row_text)
            
            text = "\n".join(paragraphs)
            logger.info(f"Extracted text from DOCX: {file_path.name}")
            return [(1, text)]
        except Exception as e:
            logger.error(f"Failed to extract text from DOCX: {file_path.name}, error: {str(e)}")
            raise

    @staticmethod
    def _extract_text_from_pptx(file_path: Path) -> List[Tuple[int, str]]:
        """Extract text from PPTX file"""
        try:
            prs = Presentation(file_path)
            slides_text = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                
                if slide_text:
                    slides_text.append((slide_num, "\n".join(slide_text)))
            
            logger.info(f"Extracted {len(slides_text)} slides from PPTX: {file_path.name}")
            return slides_text
        except Exception as e:
            logger.error(f"Failed to extract text from PPTX: {file_path.name}, error: {str(e)}")
            raise

    @staticmethod
    def _extract_text_from_xlsx(file_path: Path) -> List[Tuple[int, str]]:
        """Extract text from XLSX file"""
        try:
            wb = load_workbook(file_path, read_only=True)
            sheets_text = []
            
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                rows_text = []
                
                for row in ws.iter_rows(values_only=True):
                    row_text = "\t".join(str(cell) if cell is not None else "" for cell in row)
                    if row_text.strip():
                        rows_text.append(row_text)
                
                if rows_text:
                    sheets_text.append((1, f"Sheet: {sheet_name}\n" + "\n".join(rows_text)))
            
            wb.close()
            logger.info(f"Extracted {len(sheets_text)} sheets from XLSX: {file_path.name}")
            return sheets_text
        except Exception as e:
            logger.error(f"Failed to extract text from XLSX: {file_path.name}, error: {str(e)}")
            raise

    @staticmethod
    def _extract_text_from_csv(file_path: Path) -> List[Tuple[int, str]]:
        """Extract text from CSV file"""
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                content = f.read()
            logger.info(f"Extracted text from CSV: {file_path.name}")
            return [(1, content)]
        except Exception as e:
            logger.error(f"Failed to extract text from CSV: {file_path.name}, error: {str(e)}")
            raise

    @staticmethod
    def _extract_text_from_json(file_path: Path) -> List[Tuple[int, str]]:
        """Extract text from JSON file"""
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                content = f.read()
            logger.info(f"Extracted text from JSON: {file_path.name}")
            return [(1, content)]
        except Exception as e:
            logger.error(f"Failed to extract text from JSON: {file_path.name}, error: {str(e)}")
            raise

    @staticmethod
    def _extract_text_from_xml(file_path: Path) -> List[Tuple[int, str]]:
        """Extract text from XML file"""
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                content = f.read()
            logger.info(f"Extracted text from XML: {file_path.name}")
            return [(1, content)]
        except Exception as e:
            logger.error(f"Failed to extract text from XML: {file_path.name}, error: {str(e)}")
            raise

    @staticmethod
    def _extract_text_from_html(file_path: Path) -> List[Tuple[int, str]]:
        """Extract text from HTML file"""
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                content = f.read()
            logger.info(f"Extracted text from HTML: {file_path.name}")
            return [(1, content)]
        except Exception as e:
            logger.error(f"Failed to extract text from HTML: {file_path.name}, error: {str(e)}")
            raise

    @staticmethod
    def _extract_text_from_image(file_path: Path) -> List[Tuple[int, str]]:
        """Extract text from image file using OCR (placeholder)"""
        try:
            img = Image.open(file_path)
            img_info = f"Image: {file_path.name}\nFormat: {img.format}\nSize: {img.size}\nMode: {img.mode}"
            
            img_bytes = io.BytesIO()
            img.save(img_bytes, format=img.format or 'PNG')
            
            logger.info(f"Processed image: {file_path.name}")
            return [(1, img_info)]
        except Exception as e:
            logger.error(f"Failed to process image: {file_path.name}, error: {str(e)}")
            raise

    @staticmethod
    def _get_odf_text(element) -> str:
        """Extract text from an ODF element recursively"""
        text = ""
        if hasattr(element, 'nodeType') and element.nodeType == 3:
            text = str(element)
        elif hasattr(element, 'hasChildNodes') and element.hasChildNodes():
            child = element.firstChild
            while child is not None:
                text += DocumentProcessor._get_odf_text(child)
                child = child.nextSibling
        return text

    @staticmethod
    def _extract_text_from_odt(file_path: Path) -> List[Tuple[int, str]]:
        """Extract text from ODT file"""
        try:
            doc = load(str(file_path))
            paragraphs = []
            
            for element in doc.getElementsByType(P):
                text = DocumentProcessor._get_odf_text(element).strip()
                if text:
                    paragraphs.append(text)
            
            text = "\n".join(paragraphs)
            logger.info(f"Extracted text from ODT: {file_path.name}")
            return [(1, text)]
        except Exception as e:
            logger.error(f"Failed to extract text from ODT: {file_path.name}, error: {str(e)}")
            raise

    @staticmethod
    def _extract_text_from_ods(file_path: Path) -> List[Tuple[int, str]]:
        """Extract text from ODS file"""
        try:
            doc = load(str(file_path))
            sheets_text = []
            
            for table in doc.getElementsByType(Table):
                table_name = table.getAttribute('name') or 'Sheet'
                rows_text = []
                
                for row in table.getElementsByType(TableRow):
                    row_text = []
                    for cell in row.getElementsByType(TableCell):
                        cell_text = ""
                        for p in cell.getElementsByType(P):
                            cell_text += DocumentProcessor._get_odf_text(p) + " "
                        row_text.append(cell_text.strip())
                    
                    row_str = "\t".join(row_text)
                    if row_str.strip():
                        rows_text.append(row_str)
                
                if rows_text:
                    sheets_text.append((1, f"Sheet: {table_name}\n" + "\n".join(rows_text)))
            
            logger.info(f"Extracted {len(sheets_text)} sheets from ODS: {file_path.name}")
            return sheets_text
        except Exception as e:
            logger.error(f"Failed to extract text from ODS: {file_path.name}, error: {str(e)}")
            raise

    def process(self, file_path: Path) -> List[Tuple[int, str]]:
        """Process document file and return list of (page_number, text)"""
        ext = file_path.suffix.lower()
        
        if ext not in self.SUPPORTED_EXTENSIONS:
            raise ValueError(f"Unsupported file type: {ext}. Supported types: {self.SUPPORTED_EXTENSIONS}")
        
        logger.info(f"Processing document: {file_path.name}")
        
        if ext == ".pdf":
            return self._extract_text_from_pdf(file_path)
        elif ext in (".md", ".markdown"):
            return self._extract_text_from_markdown(file_path)
        elif ext == ".txt":
            return self._extract_text_from_txt(file_path)
        elif ext == ".docx":
            return self._extract_text_from_docx(file_path)
        elif ext == ".doc":
            return self._extract_text_from_docx(file_path)
        elif ext == ".pptx":
            return self._extract_text_from_pptx(file_path)
        elif ext == ".ppt":
            return self._extract_text_from_pptx(file_path)
        elif ext == ".xlsx":
            return self._extract_text_from_xlsx(file_path)
        elif ext == ".xls":
            return self._extract_text_from_xlsx(file_path)
        elif ext == ".csv":
            return self._extract_text_from_csv(file_path)
        elif ext == ".json":
            return self._extract_text_from_json(file_path)
        elif ext == ".xml":
            return self._extract_text_from_xml(file_path)
        elif ext in (".html", ".htm"):
            return self._extract_text_from_html(file_path)
        elif ext in (".jpg", ".jpeg", ".png", ".gif", ".bmp"):
            return self._extract_text_from_image(file_path)
        elif ext == ".odt":
            return self._extract_text_from_odt(file_path)
        elif ext == ".ods":
            return self._extract_text_from_ods(file_path)
        else:
            raise ValueError(f"Unsupported file type: {ext}")